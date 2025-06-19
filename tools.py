import wikipedia
from io import BytesIO
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
import logging
logging.basicConfig(level=logging.INFO)

import requests
import json
from typing import List, Dict, Optional, Union
from urllib.parse import quote
import re
from dataclasses import dataclass


@dataclass
class WikipediaPage:
    """Data class to represent a Wikipedia page"""
    title: str
    content: str
    url: str
    summary: str
    categories: List[str]
    links: List[str]


class WikipediaAgent:
    """
    A comprehensive Wikipedia tool for AI agents.
    Provides search, retrieval, and content processing capabilities.
    """
    
    def __init__(self, language: str = "en", user_agent: str = "WikipediaAgent/1.0"):
        """
        Initialize the Wikipedia agent.
        
        Args:
            language: Wikipedia language code (default: "en")
            user_agent: User agent string for API requests
        """
        self.language = language
        self.base_url = f"https://{language}.wikipedia.org/api/rest_v1"
        self.api_url = f"https://{language}.wikipedia.org/w/api.php"
        self.headers = {
            "User-Agent": user_agent,
            "Accept": "application/json"
        }
    
    def search(self, query: str, limit: int = 10) -> List[Dict[str, Union[str, int]]]:
        """
        Search Wikipedia for pages matching the query.
        
        Args:
            query: Search query string
            limit: Maximum number of results to return
            
        Returns:
            List of dictionaries containing search results
        """
        params = {
            "action": "query",
            "format": "json",
            "list": "search",
            "srsearch": query,
            "srlimit": limit,
            "srprop": "snippet|titlesnippet|size|timestamp"
        }
        
        try:
            response = requests.get(self.api_url, params=params, headers=self.headers)
            response.raise_for_status()
            data = response.json()
            
            results = []
            for item in data.get("query", {}).get("search", []):
                results.append({
                    "title": item["title"],
                    "snippet": self._clean_html(item.get("snippet", "")),
                    "size": item.get("size", 0),
                    "timestamp": item.get("timestamp", ""),
                    "url": f"https://{self.language}.wikipedia.org/wiki/{quote(item['title'])}"
                })
            
            return results
            
        except requests.RequestException as e:
            raise Exception(f"Search failed: {str(e)}")
    
    def get_page(self, title: str, include_links: bool = True) -> Optional[WikipediaPage]:
        """
        Retrieve a complete Wikipedia page.
        
        Args:
            title: Page title
            include_links: Whether to extract internal links
            
        Returns:
            WikipediaPage object or None if not found
        """
        try:
            # Get page content
            content_url = f"{self.base_url}/page/html/{quote(title)}"
            content_response = requests.get(content_url, headers=self.headers)
            
            if content_response.status_code == 404:
                return None
            
            content_response.raise_for_status()
            html_content = content_response.text
            
            # Get page summary
            summary_url = f"{self.base_url}/page/summary/{quote(title)}"
            summary_response = requests.get(summary_url, headers=self.headers)
            summary_data = summary_response.json() if summary_response.status_code == 200 else {}
            
            # Extract plain text content
            plain_text = self._extract_text_from_html(html_content)
            
            # Get categories and links
            categories = self._get_page_categories(title)
            links = self._get_page_links(title) if include_links else []
            
            return WikipediaPage(
                title=title,
                content=plain_text,
                url=f"https://{self.language}.wikipedia.org/wiki/{quote(title)}",
                summary=summary_data.get("extract", ""),
                categories=categories,
                links=links
            )
            
        except requests.RequestException as e:
            raise Exception(f"Failed to retrieve page '{title}': {str(e)}")
    
    def get_summary(self, title: str, sentences: int = 3) -> Optional[str]:
        """
        Get a summary of a Wikipedia page.
        
        Args:
            title: Page title
            sentences: Number of sentences in summary
            
        Returns:
            Summary text or None if not found
        """
        try:
            params = {
                "action": "query",
                "format": "json",
                "prop": "extracts",
                "exintro": True,
                "explaintext": True,
                "exsentences": sentences,
                "titles": title
            }
            
            response = requests.get(self.api_url, params=params, headers=self.headers)
            response.raise_for_status()
            data = response.json()
            
            pages = data.get("query", {}).get("pages", {})
            for page_id, page_data in pages.items():
                if page_id != "-1":  # Page exists
                    return page_data.get("extract", "")
            
            return None
            
        except requests.RequestException as e:
            raise Exception(f"Failed to get summary for '{title}': {str(e)}")
    
    def get_random_page(self) -> Optional[str]:
        """
        Get a random Wikipedia page title.
        
        Returns:
            Random page title or None if failed
        """
        try:
            params = {
                "action": "query",
                "format": "json",
                "list": "random",
                "rnnamespace": 0,
                "rnlimit": 1
            }
            
            response = requests.get(self.api_url, params=params, headers=self.headers)
            response.raise_for_status()
            data = response.json()
            
            random_pages = data.get("query", {}).get("random", [])
            if random_pages:
                return random_pages[0]["title"]
            
            return None
            
        except requests.RequestException as e:
            raise Exception(f"Failed to get random page: {str(e)}")
    
    def get_page_links(self, title: str, limit: int = 50) -> List[str]:
        """
        Get internal links from a Wikipedia page.
        
        Args:
            title: Page title
            limit: Maximum number of links to return
            
        Returns:
            List of linked page titles
        """
        return self._get_page_links(title, limit)
    
    def get_page_categories(self, title: str) -> List[str]:
        """
        Get categories for a Wikipedia page.
        
        Args:
            title: Page title
            
        Returns:
            List of category names
        """
        return self._get_page_categories(title)
    
    def _get_page_links(self, title: str, limit: int = 50) -> List[str]:
        """Internal method to get page links"""
        try:
            params = {
                "action": "query",
                "format": "json",
                "prop": "links",
                "titles": title,
                "pllimit": limit,
                "plnamespace": 0
            }
            
            response = requests.get(self.api_url, params=params, headers=self.headers)
            response.raise_for_status()
            data = response.json()
            
            links = []
            pages = data.get("query", {}).get("pages", {})
            for page_data in pages.values():
                if "links" in page_data:
                    links.extend([link["title"] for link in page_data["links"]])
            
            return links
            
        except requests.RequestException:
            return []
    
    def _get_page_categories(self, title: str) -> List[str]:
        """Internal method to get page categories"""
        try:
            params = {
                "action": "query",
                "format": "json",
                "prop": "categories",
                "titles": title,
                "cllimit": 50
            }
            
            response = requests.get(self.api_url, params=params, headers=self.headers)
            response.raise_for_status()
            data = response.json()
            
            categories = []
            pages = data.get("query", {}).get("pages", {})
            for page_data in pages.values():
                if "categories" in page_data:
                    categories.extend([
                        cat["title"].replace("Category:", "") 
                        for cat in page_data["categories"]
                    ])
            
            return categories
            
        except requests.RequestException:
            return []
    
    def _extract_text_from_html(self, html: str) -> str:
        """Extract plain text from HTML content"""
        # Remove HTML tags
        text = re.sub(r'<[^>]+>', '', html)
        # Clean up whitespace
        text = re.sub(r'\s+', ' ', text)
        # Remove special characters and references
        text = re.sub(r'\[\d+\]', '', text)  # Remove reference numbers
        return text.strip()
    
    def _clean_html(self, text: str) -> str:
        """Clean HTML entities and tags from text"""
        # Remove HTML tags
        text = re.sub(r'<[^>]+>', '', text)
        # Replace common HTML entities
        text = text.replace('&quot;', '"').replace('&amp;', '&').replace('&lt;', '<').replace('&gt;', '>')
        return text.strip()


# Example usage and utility functions
def create_knowledge_base(agent: WikipediaAgent, topics: List[str]) -> Dict[str, WikipediaPage]:
    """
    Create a knowledge base from a list of topics.
    
    Args:
        agent: WikipediaAgent instance
        topics: List of topic titles
        
    Returns:
        Dictionary mapping topics to WikipediaPage objects
    """
    knowledge_base = {}
    
    for topic in topics:
        try:
            page = agent.get_page(topic)
            if page:
                knowledge_base[topic] = page
                print(f"Added: {topic}")
            else:
                print(f"Not found: {topic}")
        except Exception as e:
            print(f"Error with {topic}: {str(e)}")
    
    return knowledge_base


def search_and_summarize(agent: WikipediaAgent, query: str, num_results: int = 3) -> List[Dict]:
    """
    Search for topics and get summaries.
    
    Args:
        agent: WikipediaAgent instance
        query: Search query
        num_results: Number of results to summarize
        
    Returns:
        List of dictionaries with titles and summaries
    """
    search_results = agent.search(query, limit=num_results)
    summaries = []
    
    for result in search_results:
        summary = agent.get_summary(result["title"])
        summaries.append({
            "title": result["title"],
            "url": result["url"],
            "summary": summary,
            "snippet": result["snippet"]
        })
    
    return summaries


# Example usage
if __name__ == "__main__":
    # Initialize the agent
    wiki = WikipediaAgent()
    
    # Search for pages
    results = wiki.search("artificial intelligence", limit=5)
    print("Search Results:")
    for result in results:
        print(f"- {result['title']}: {result['snippet'][:100]}...")
    
    # Get a specific page
    page = wiki.get_page("Machine learning")
    if page:
        print(f"\nPage: {page.title}")
        print(f"Summary: {page.summary[:200]}...")
        print(f"Categories: {page.categories[:5]}")
        print(f"Links: {page.links[:10]}")
    
    # Get summary
    summary = wiki.get_summary("Deep learning", sentences=2)
    print(f"\nDeep Learning Summary: {summary}")
    
    # Create knowledge base
    topics = ["Machine learning", "Neural network", "Natural language processing"]
    kb = create_knowledge_base(wiki, topics)
    print(f"\nKnowledge base created with {len(kb)} entries") 
    """
PowerPoint Tools using python-pptx
Install required library: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

class PowerPointTool:
    def __init__(self):
        self.prs = Presentation()
        
    def create_new_presentation(self):
        """Create a new blank presentation"""
        self.prs = Presentation()
        return self
    
    def load_presentation(self, file_path):
        """Load an existing presentation"""
        if os.path.exists(file_path):
            self.prs = Presentation(file_path)
            print(f"Loaded presentation: {file_path}")
        else:
            print(f"File not found: {file_path}")
        return self
    
    def add_title_slide(self, title, subtitle=""):
        """Add a title slide"""
        slide_layout = self.prs.slide_layouts[0]  # Title slide layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        subtitle_shape.text = subtitle
        
        return slide
    
    def add_content_slide(self, title, content_points):
        """Add a slide with title and bullet points"""
        slide_layout = self.prs.slide_layouts[1]  # Title and content layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]
        
        title_shape.text = title
        
        text_frame = content_shape.text_frame
        text_frame.clear()  # Clear existing text
        
        for i, point in enumerate(content_points):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = point
            p.level = 0  # Bullet level
        
        return slide
    
    def add_blank_slide(self):
        """Add a blank slide"""
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        return slide
    
    def add_text_box(self, slide, text, left, top, width, height):
        """Add a text box to a slide"""
        text_box = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        text_frame = text_box.text_frame
        text_frame.text = text
        return text_box
    
    def add_image(self, slide, image_path, left, top, width=None, height=None):
        """Add an image to a slide"""
        if not os.path.exists(image_path):
            print(f"Image not found: {image_path}")
            return None
            
        if width and height:
            pic = slide.shapes.add_picture(
                image_path, Inches(left), Inches(top), 
                Inches(width), Inches(height)
            )
        else:
            pic = slide.shapes.add_picture(
                image_path, Inches(left), Inches(top)
            )
        return pic
    
    def format_text(self, text_frame, font_name="Arial", font_size=18, 
                   bold=False, italic=False, color=None):
        """Format text in a text frame"""
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                font = run.font
                font.name = font_name
                font.size = Pt(font_size)
                font.bold = bold
                font.italic = italic
                if color:
                    font.color.rgb = RGBColor(*color)  # RGB tuple
    
    def set_slide_background_color(self, slide, rgb_color):
        """Set slide background color"""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*rgb_color)
    
    def get_slide_count(self):
        """Get number of slides"""
        return len(self.prs.slides)
    
    def get_slide_text(self, slide_index):
        """Extract all text from a specific slide"""
        if slide_index >= len(self.prs.slides):
            return "Slide index out of range"
        
        slide = self.prs.slides[slide_index]
        text_content = []
        
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_content.append(shape.text)
        
        return "\n".join(text_content)
    
    def extract_all_text(self):
        """Extract text from all slides"""
        all_text = {}
        for i, slide in enumerate(self.prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            all_text[f"Slide {i+1}"] = slide_text
        return all_text
    
    def save_presentation(self, filename):
        """Save the presentation"""
        self.prs.save(filename)
        print(f"Presentation saved as: {filename}")
    
    def create_sample_presentation(self):
        """Create a sample presentation with different slide types"""
        # Title slide
        self.add_title_slide("Sample Presentation", "Created with Python")
        
        # Content slide
        content_points = [
            "First bullet point",
            "Second bullet point", 
            "Third bullet point"
        ]
        self.add_content_slide("Main Topics", content_points)
        
        # Blank slide with custom content
        blank_slide = self.add_blank_slide()
        self.add_text_box(blank_slide, "Custom Text Box", 1, 1, 8, 2)
        
        return self

# Example usage functions
def create_basic_presentation():
    """Example: Create a basic presentation"""
    ppt_tool = PowerPointTool()
    
    # Create title slide
    ppt_tool.add_title_slide("My Presentation", "Subtitle here")
    
    # Add content slides
    topics = ["Introduction", "Main Content", "Conclusion"]
    for topic in topics:
        points = [f"Point 1 about {topic}", f"Point 2 about {topic}"]
        ppt_tool.add_content_slide(topic, points)
    
    # Save presentation
    ppt_tool.save_presentation("basic_presentation.pptx")
    
    return ppt_tool

def analyze_existing_presentation(file_path):
    """Example: Analyze an existing presentation"""
    ppt_tool = PowerPointTool()
    ppt_tool.load_presentation(file_path)
    
    # Get basic info
    print(f"Number of slides: {ppt_tool.get_slide_count()}")
    
    # Extract all text
    all_text = ppt_tool.extract_all_text()
    for slide_name, text_list in all_text.items():
        print(f"\n{slide_name}:")
        for text in text_list:
            print(f"  - {text}")
    
    return all_text

def create_formatted_presentation():
    """Example: Create a presentation with custom formatting"""
    ppt_tool = PowerPointTool()
    
    # Title slide
    title_slide = ppt_tool.add_title_slide("Formatted Presentation")
    
    # Add a content slide
    content_slide = ppt_tool.add_content_slide("Features", [
        "Custom formatting",
        "Colors and fonts",
        "Images and text boxes"
    ])
    
    # Add formatted text box to blank slide
    blank_slide = ppt_tool.add_blank_slide()
    text_box = ppt_tool.add_text_box(blank_slide, "Formatted Text", 1, 1, 6, 2)
    
    # Format the text (red color)
    ppt_tool.format_text(text_box.text_frame, font_size=24, bold=True, 
                        color=(255, 0, 0))
    
    # Set background color (light blue)
    ppt_tool.set_slide_background_color(blank_slide, (173, 216, 230))
    
    ppt_tool.save_presentation("formatted_presentation.pptx")
    return ppt_tool

# Main execution
if __name__ == "__main__":
    # Create a sample presentation
    print("Creating sample presentation...")
    ppt_tool = PowerPointTool()
    ppt_tool.create_sample_presentation()
    ppt_tool.save_presentation("sample_presentation.pptx")
    
    # Create a basic presentation
    print("\nCreating basic presentation...")
    create_basic_presentation()
    
    # Create formatted presentation
    print("\nCreating formatted presentation...")
    create_formatted_presentation()
    
    print("\nAll presentations created successfully!")
